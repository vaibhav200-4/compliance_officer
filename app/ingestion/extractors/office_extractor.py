from pathlib import Path

from docx import Document as DocxDocument
from pptx import Presentation
from langchain_core.documents import Document


class OfficeExtractor:

    def extract(self, file_path: Path):

        suffix = file_path.suffix.lower()

        if suffix == ".docx":

            doc = DocxDocument(file_path)

            text = "\n".join(
                p.text
                for p in doc.paragraphs
                if p.text.strip()
            )

        elif suffix == ".pptx":

            prs = Presentation(file_path)

            slides = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides.append(shape.text)

            text = "\n".join(slides)

        else:
            raise ValueError("Unsupported Office document.")

        return [
            Document(
                page_content=text,
                metadata={
                    "source": file_path.name,
                    "page": 1,
                },
            )
        ]